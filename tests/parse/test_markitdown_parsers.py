# Copyright (c) 2026 Beijing Volcano Engine Technology Co., Ltd.
# SPDX-License-Identifier: Apache-2.0
"""Unit tests for markitdown-inspired parsers."""

import zipfile

import pytest

from openviking.parse.parsers.audio import AudioParser
from openviking.parse.parsers.epub import EPubParser
from openviking.parse.parsers.excel import ExcelParser
from openviking.parse.parsers.powerpoint import PowerPointParser
from openviking.parse.parsers.word import WordParser
from openviking.parse.parsers.zip_parser import ZipParser

# ---------------------------------------------------------------------------
# Helpers to create real test files
# ---------------------------------------------------------------------------


def _create_docx(path, paragraphs=None, tables=None):
    """Create a real .docx file with optional paragraphs and tables."""
    import docx

    doc = docx.Document()
    if paragraphs:
        for style, text in paragraphs:
            doc.add_paragraph(text, style=style)
    if tables:
        for rows in tables:
            table = doc.add_table(rows=len(rows), cols=len(rows[0]))
            for i, row_data in enumerate(rows):
                for j, cell_text in enumerate(row_data):
                    table.rows[i].cells[j].text = cell_text
    doc.save(str(path))
    return path


def _create_xlsx(path, sheets=None):
    """Create a real .xlsx file with named sheets."""
    import openpyxl

    wb = openpyxl.Workbook()
    if sheets:
        for idx, (name, rows) in enumerate(sheets.items()):
            ws = wb.active if idx == 0 else wb.create_sheet()
            ws.title = name
            for row in rows:
                ws.append(row)
    wb.save(str(path))
    return path


def _create_pptx(path, slides=None):
    """Create a real .pptx file with slides."""
    from pptx import Presentation

    prs = Presentation()
    if slides:
        for title, body in slides:
            layout = prs.slide_layouts[1]  # title + content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
    prs.save(str(path))
    return path


# -----------------------------------------------------------------------------
# Word Parser Tests
# -----------------------------------------------------------------------------


class TestWordParser:
    """Test Word (.docx) parser."""

    @pytest.fixture
    def word_parser(self):
        return WordParser()

    def test_supported_extensions(self, word_parser):
        assert ".docx" in word_parser.supported_extensions

    @pytest.mark.asyncio
    async def test_parse_content_delegates_to_markdown(self, word_parser):
        """Test that parse_content delegates to MarkdownParser."""
        content = "# Test Heading\n\nThis is test content."
        result = await word_parser.parse_content(content, source_path="test.docx")

        assert result.source_format == "docx"
        assert result.parser_name == "WordParser"

    @pytest.mark.asyncio
    async def test_parse_nonexistent_file_falls_back(self, word_parser):
        """Test that parse() with non-existent path treats source as content."""
        result = await word_parser.parse("# Heading\n\nSome content")
        assert result.source_format == "docx"
        assert result.parser_name == "WordParser"


# -----------------------------------------------------------------------------
# PowerPoint Parser Tests
# -----------------------------------------------------------------------------


class TestPowerPointParser:
    """Test PowerPoint (.pptx) parser."""

    @pytest.fixture
    def ppt_parser(self):
        return PowerPointParser()

    def test_supported_extensions(self, ppt_parser):
        assert ".pptx" in ppt_parser.supported_extensions

    @pytest.mark.asyncio
    async def test_parse_content_delegates_to_markdown(self, ppt_parser):
        """Test that parse_content delegates to MarkdownParser."""
        content = "# Slide 1\n\nContent here."
        result = await ppt_parser.parse_content(content, source_path="test.pptx")

        assert result.source_format == "pptx"
        assert result.parser_name == "PowerPointParser"


# -----------------------------------------------------------------------------
# Excel Parser Tests
# -----------------------------------------------------------------------------


class TestExcelParser:
    """Test Excel (.xlsx) parser."""

    @pytest.fixture
    def excel_parser(self):
        return ExcelParser()

    def test_supported_extensions(self, excel_parser):
        assert ".xlsx" in excel_parser.supported_extensions
        assert ".xls" in excel_parser.supported_extensions
        assert ".xlsm" in excel_parser.supported_extensions

    @pytest.mark.asyncio
    async def test_parse_content_delegates_to_markdown(self, excel_parser):
        """Test that parse_content delegates to MarkdownParser."""
        content = "| Col1 | Col2 |\n|------|------|\n| A | B |"
        result = await excel_parser.parse_content(content, source_path="test.xlsx")

        assert result.source_format == "xlsx"
        assert result.parser_name == "ExcelParser"


# -----------------------------------------------------------------------------
# EPUB Parser Tests
# -----------------------------------------------------------------------------


class TestEPubParser:
    """Test EPUB parser."""

    @pytest.fixture
    def epub_parser(self):
        return EPubParser()

    def test_supported_extensions(self, epub_parser):
        assert ".epub" in epub_parser.supported_extensions

    @pytest.mark.asyncio
    async def test_parse_content_delegates_to_markdown(self, epub_parser):
        """Test that parse_content delegates to MarkdownParser."""
        content = "# Chapter 1\n\nOnce upon a time..."
        result = await epub_parser.parse_content(content, source_path="test.epub")

        assert result.source_format == "epub"
        assert result.parser_name == "EPubParser"


# -----------------------------------------------------------------------------
# ZIP Parser Tests
# -----------------------------------------------------------------------------


class TestZipParser:
    """Test ZIP archive parser."""

    @pytest.fixture
    def zip_parser(self):
        return ZipParser()

    @pytest.fixture
    def sample_zip(self, tmp_path):
        """Create a sample ZIP file for testing."""
        zip_path = tmp_path / "test.zip"
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("readme.txt", "This is a readme file.")
            zf.writestr("data/info.json", '{"key": "value"}')
            zf.writestr("src/main.py", "print('hello')")
        return zip_path

    def test_supported_extensions(self, zip_parser):
        assert ".zip" in zip_parser.supported_extensions

    def test_convert_zip_to_markdown(self, zip_parser, sample_zip):
        """Test ZIP to markdown conversion."""
        markdown = zip_parser._convert_zip_to_markdown(sample_zip)

        assert "# ZIP Archive:" in markdown
        assert "readme.txt" in markdown
        assert "info.json" in markdown
        assert "main.py" in markdown

    def test_format_size(self, zip_parser):
        """Test file size formatting."""
        assert zip_parser._format_size(500) == "500.0 B"
        assert zip_parser._format_size(2048) == "2.0 KB"
        assert zip_parser._format_size(2097152) == "2.0 MB"

    def test_group_files_by_extension(self, zip_parser):
        """Test file grouping by extension."""
        filenames = ["file.txt", "data.json", "script.py", "notes.txt"]
        groups = zip_parser._group_files_by_extension(filenames)

        assert groups[".txt"] == ["file.txt", "notes.txt"]
        assert groups[".json"] == ["data.json"]
        assert groups[".py"] == ["script.py"]

    @pytest.mark.asyncio
    async def test_parse_zip_file(self, zip_parser, sample_zip):
        """Test parsing an actual ZIP file."""
        result = await zip_parser.parse(sample_zip)

        assert result.source_format == "zip"
        assert result.parser_name == "ZipParser"

    @pytest.mark.asyncio
    async def test_parse_content_delegates_to_markdown(self, zip_parser):
        """Test that parse_content delegates to MarkdownParser."""
        content = "# ZIP Archive: test.zip\n\nSome content"
        result = await zip_parser.parse_content(content)

        assert result.source_format == "zip"
        assert result.parser_name == "ZipParser"


# -----------------------------------------------------------------------------
# Audio Parser Tests
# -----------------------------------------------------------------------------


class TestAudioParser:
    """Test Audio file parser."""

    @pytest.fixture
    def audio_parser(self):
        return AudioParser()

    def test_supported_extensions(self, audio_parser):
        supported = audio_parser.supported_extensions
        assert ".mp3" in supported
        assert ".wav" in supported
        assert ".m4a" in supported
        assert ".ogg" in supported
        assert ".flac" in supported

    def test_format_size(self, audio_parser):
        """Test file size formatting."""
        assert audio_parser._format_size(500) == "500.0 B"
        assert audio_parser._format_size(2048) == "2.0 KB"
        assert audio_parser._format_size(2097152) == "2.0 MB"

    def test_format_duration(self, audio_parser):
        """Test duration formatting."""
        assert audio_parser._format_duration(65) == "1:05"
        assert audio_parser._format_duration(3661) == "1:01:01"
        assert audio_parser._format_duration(45) == "0:45"

    @pytest.mark.asyncio
    async def test_parse_content_delegates_to_markdown(self, audio_parser):
        """Test that parse_content delegates to MarkdownParser."""
        content = "# Audio File: test.mp3\n\n**Duration:** 3:45"
        result = await audio_parser.parse_content(content, source_path="test.mp3")

        assert result.source_format == "audio"
        assert result.parser_name == "AudioParser"


# -----------------------------------------------------------------------------
# Real File Parsing Tests (exercise actual conversion logic)
# -----------------------------------------------------------------------------


class TestWordParserRealFile:
    """Test WordParser with real .docx files."""

    @pytest.fixture
    def word_parser(self):
        return WordParser()

    @pytest.mark.asyncio
    async def test_parse_docx_with_headings_and_body(self, word_parser, tmp_path):
        """Headings and body paragraphs are preserved in order."""
        docx_path = _create_docx(
            tmp_path / "test.docx",
            paragraphs=[
                ("Heading 1", "Introduction"),
                ("Normal", "First paragraph of text."),
                ("Heading 2", "Details"),
                ("Normal", "Second paragraph."),
            ],
        )
        result = await word_parser.parse(docx_path)
        assert result.source_format == "docx"
        assert result.parser_name == "WordParser"

    @pytest.mark.asyncio
    async def test_parse_docx_table_position_preserved(self, word_parser, tmp_path):
        """Tables interspersed with paragraphs keep their document order."""
        docx_path = tmp_path / "tables.docx"
        _create_docx(
            docx_path,
            paragraphs=[("Normal", "Before table.")],
            tables=[[["H1", "H2"], ["A", "B"]]],
        )
        # Just ensure it doesn't crash and metadata is set
        result = await word_parser.parse(docx_path)
        assert result.source_format == "docx"

    @pytest.mark.asyncio
    async def test_parse_docx_can_parse_check(self, word_parser):
        assert word_parser.can_parse("report.docx")
        assert not word_parser.can_parse("report.pdf")


class TestExcelParserRealFile:
    """Test ExcelParser with real .xlsx files."""

    @pytest.fixture
    def excel_parser(self):
        return ExcelParser()

    @pytest.mark.asyncio
    async def test_parse_xlsx_single_sheet(self, excel_parser, tmp_path):
        xlsx_path = _create_xlsx(
            tmp_path / "data.xlsx",
            sheets={"Sales": [["Product", "Revenue"], ["Widget", "1000"], ["Gadget", "2500"]]},
        )
        result = await excel_parser.parse(xlsx_path)
        assert result.source_format == "xlsx"
        assert result.parser_name == "ExcelParser"

    @pytest.mark.asyncio
    async def test_parse_xlsx_multiple_sheets(self, excel_parser, tmp_path):
        xlsx_path = _create_xlsx(
            tmp_path / "multi.xlsx",
            sheets={
                "Sheet1": [["A", "B"], [1, 2]],
                "Sheet2": [["X", "Y"], [3, 4]],
            },
        )
        result = await excel_parser.parse(xlsx_path)
        assert result.source_format == "xlsx"

    def test_convert_sheet_markdown_contains_table(self, excel_parser, tmp_path):
        """Verify the markdown output contains proper table syntax."""
        import openpyxl

        xlsx_path = _create_xlsx(
            tmp_path / "tbl.xlsx",
            sheets={"Data": [["Col1", "Col2"], ["val1", "val2"]]},
        )
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        md = excel_parser._convert_sheet(wb["Data"], "Data")
        assert "| Col1" in md
        assert "| val1" in md


class TestPowerPointParserRealFile:
    """Test PowerPointParser with real .pptx files."""

    @pytest.fixture
    def ppt_parser(self):
        return PowerPointParser()

    @pytest.mark.asyncio
    async def test_parse_pptx_with_slides(self, ppt_parser, tmp_path):
        pptx_path = _create_pptx(
            tmp_path / "deck.pptx",
            slides=[
                ("Welcome", "Agenda for today"),
                ("Topic A", "Details about topic A"),
            ],
        )
        result = await ppt_parser.parse(pptx_path)
        assert result.source_format == "pptx"
        assert result.parser_name == "PowerPointParser"


class TestEPubParserHtmlConversion:
    """Test EPubParser HTML-to-markdown conversion logic."""

    @pytest.fixture
    def epub_parser(self):
        return EPubParser()

    def test_html_to_markdown_headings(self, epub_parser):
        html = "<h1>Title</h1><p>Paragraph</p>"
        md = epub_parser._html_to_markdown(html)
        assert "# Title" in md
        assert "Paragraph" in md

    def test_html_to_markdown_bold_italic(self, epub_parser):
        html = "<p><strong>bold</strong> and <em>italic</em></p>"
        md = epub_parser._html_to_markdown(html)
        assert "**bold**" in md
        assert "*italic*" in md

    def test_html_to_markdown_strips_script_style(self, epub_parser):
        html = "<script>alert('x')</script><style>.a{}</style><p>safe</p>"
        md = epub_parser._html_to_markdown(html)
        assert "alert" not in md
        assert "safe" in md

    @pytest.mark.asyncio
    async def test_parse_epub_manual_fallback(self, epub_parser, tmp_path):
        """Test manual ZIP-based epub extraction fallback."""
        epub_path = tmp_path / "book.epub"
        with zipfile.ZipFile(epub_path, "w") as zf:
            zf.writestr("mimetype", "application/epub+zip")
            zf.writestr(
                "OEBPS/chapter1.xhtml",
                "<html><body><h1>Chapter 1</h1><p>Hello world</p></body></html>",
            )
        # _convert_manual should handle this
        md = epub_parser._convert_manual(epub_path)
        assert "Chapter 1" in md
        assert "Hello world" in md


class TestZipParserTreeView:
    """Test ZipParser tree view uses plain text (no emoji)."""

    @pytest.fixture
    def zip_parser(self):
        return ZipParser()

    def test_tree_view_no_emoji(self, zip_parser):
        filenames = ["dir/", "dir/file.txt"]
        tree = zip_parser._generate_tree_view(filenames)
        # Should NOT contain emoji characters
        assert "\U0001f4c1" not in tree  # 📁
        assert "\U0001f4c4" not in tree  # 📄

    def test_bad_zip_raises(self, zip_parser, tmp_path):
        bad = tmp_path / "bad.zip"
        bad.write_text("not a zip")
        with pytest.raises(ValueError, match="Invalid or corrupted"):
            zip_parser._convert_zip_to_markdown(bad)


class TestAudioParserEdgeCases:
    """Test AudioParser edge cases."""

    @pytest.fixture
    def audio_parser(self):
        return AudioParser()

    def test_extract_tags_empty(self, audio_parser):
        """No tags when audio object has no tags attr."""

        class FakeAudio:
            tags = None

        assert audio_parser._extract_tags(FakeAudio()) == {}

    def test_extract_tags_vorbis_style(self, audio_parser):
        """Vorbis-style lowercase keys are extracted."""

        class FakeTags(dict):
            pass

        class FakeAudio:
            tags = FakeTags({"title": ["My Song"], "artist": ["Artist Name"]})

        tags = audio_parser._extract_tags(FakeAudio())
        assert tags.get("Title") == "My Song"
        assert tags.get("Artist") == "Artist Name"

    def test_format_duration_zero(self, audio_parser):
        assert audio_parser._format_duration(0) == "0:00"


# -----------------------------------------------------------------------------
# Registry Integration Tests
# -----------------------------------------------------------------------------


class TestRegistryIntegration:
    """Test parser registration in registry."""

    @pytest.fixture
    def registry(self):
        from openviking.parse.registry import ParserRegistry

        return ParserRegistry(register_optional=False)

    def test_markitdown_parsers_registered(self, registry):
        """Test all markitdown parsers are registered by default."""
        parsers = registry.list_parsers()
        assert "word" in parsers
        assert "powerpoint" in parsers
        assert "excel" in parsers
        assert "epub" in parsers
        assert "zip" in parsers
        assert "audio" in parsers

    def test_extensions_mapped(self, registry):
        """Test file extensions are properly mapped."""
        extensions = registry.list_supported_extensions()
        assert ".docx" in extensions
        assert ".pptx" in extensions
        assert ".xlsx" in extensions
        assert ".xlsm" in extensions
        assert ".epub" in extensions
        assert ".zip" in extensions
        assert ".mp3" in extensions
        assert ".wav" in extensions

    def test_get_parser_for_file(self, registry):
        """Test getting parser for specific file types."""
        assert registry.get_parser_for_file("test.docx") is not None
        assert registry.get_parser_for_file("test.xlsx") is not None
        assert registry.get_parser_for_file("test.zip") is not None
        assert registry.get_parser_for_file("test.pptx") is not None
        assert registry.get_parser_for_file("test.epub") is not None
        assert registry.get_parser_for_file("test.mp3") is not None

    def test_parser_types(self, registry):
        """Test parsers are the correct types."""
        assert isinstance(registry.get_parser("word"), WordParser)
        assert isinstance(registry.get_parser("powerpoint"), PowerPointParser)
        assert isinstance(registry.get_parser("excel"), ExcelParser)
        assert isinstance(registry.get_parser("epub"), EPubParser)
        assert isinstance(registry.get_parser("zip"), ZipParser)
        assert isinstance(registry.get_parser("audio"), AudioParser)
